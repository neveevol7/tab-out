from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml

def apply_gradient_background(pptx_path, output_path):
    prs = Presentation(pptx_path)
    
    # Colors extracted
    # Top-Left: #FFFFFF
    # Bottom-Right: #E5FDA9
    
    color1 = "FFFFFF"
    color2 = "E5FDA9"
    
    for slide in prs.slides:
        # Add a rectangle that covers the whole slide
        width = prs.slide_width
        height = prs.slide_height
        
        # We need to send this rectangle to the back
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, width, height
        )
        
        # Remove border
        shape.line.fill.background()
        
        # Apply gradient via XML (python-pptx doesn't have a high-level API for gradients)
        # This is a bit complex. Let's try a simpler approach first: solid color or image.
        # But the user asked for the gradient.
        
        fill = shape.fill
        # Accessing the underlying XML to set gradient
        fill_xml = fill._fill
        
        # Constructing the gradient XML
        # gsLst: gradient stop list
        # lin: linear gradient (angle 45 degrees is approx 2700000 in 60000ths of a degree)
        # Here we use a linear gradient from TL to BR.
        
        gradient_xml = f'''
        <a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
            <a:gsLst>
                <a:gs pos="0">
                    <a:srgbClr val="{color1}"/>
                </a:gs>
                <a:gs pos="100000">
                    <a:srgbClr val="{color2}"/>
                </a:gs>
            </a:gsLst>
            <a:lin ang="2700000" scaled="1"/>
        </a:gradFill>
        '''
        
        # Replace the fill with our gradient
        new_fill = parse_xml(gradient_xml)
        fill_xml.getparent().replace(fill_xml, new_fill)
        
        # Move to back
        # slide.shapes is a list-like. Moving to back means making it the first element in XML.
        # But in python-pptx, we can use shape.z_order or manipulate the XML.
        # Simple way: slide.shapes._spTree.insert(2, shape._element) 
        # (index 0 is nvSpPr, 1 is grpSpPr, then shapes)
        slide.shapes._spTree.remove(shape._element)
        slide.shapes._spTree.insert(2, shape._element)

    prs.save(output_path)
    print(f"Saved modified PPTX to {output_path}")

if __name__ == "__main__":
    apply_gradient_background("../Documents/ima/高校教师实训营.pptx", "../Documents/ima/高校教师实训营_new_bg.pptx")
