from PIL import Image
from pptx import Presentation
from pptx.util import Inches

def create_gradient_image(color1, color2, width, height, output_path):
    # color1, color2 are hex strings like "FFFFFF"
    c1 = tuple(int(color1[i:i+2], 16) for i in (0, 2, 4))
    c2 = tuple(int(color2[i:i+2], 16) for i in (0, 2, 4))
    
    base = Image.new('RGB', (width, height), c1)
    top = Image.new('RGB', (width, height), c2)
    
    mask = Image.new('L', (width, height))
    for y in range(height):
        for x in range(width):
            # Diagonal gradient: x+y
            # Normalize to 0-255
            # We want (0,0) to be c1 and (width, height) to be c2
            val = int(255 * (x + y) / (width + height))
            mask.putpixel((x, y), val)
            
    result = Image.composite(top, base, mask)
    result.save(output_path)
    print(f"Gradient image saved to {output_path}")

def apply_image_background(pptx_path, image_path, output_path):
    prs = Presentation(pptx_path)
    
    for slide in prs.slides:
        # Add the image
        slide.shapes.add_picture(image_path, 0, 0, prs.slide_width, prs.slide_height)
        
        # Send to back
        # The last added shape is the last element in spTree
        pic = slide.shapes[-1]
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)
        
    prs.save(output_path)
    print(f"Saved modified PPTX to {output_path}")

if __name__ == "__main__":
    color1 = "FFFFFF"
    color2 = "E5FDA9"
    # Create a 1920x1080 gradient image
    create_gradient_image(color1, color2, 1920, 1080, "gradient_bg.png")
    apply_image_background("../Documents/ima/高校教师实训营.pptx", "gradient_bg.png", "../Documents/ima/高校教师实训营_new_bg.pptx")
