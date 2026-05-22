from pptx import Presentation

def apply_image_background(pptx_path, image_path, output_path):
    prs = Presentation(pptx_path)
    
    for slide in prs.slides:
        # Add the image and scale it to cover the whole slide
        # Using a slightly larger scale to avoid any thin borders
        pic = slide.shapes.add_picture(image_path, 0, 0, prs.slide_width, prs.slide_height)
        
        # Send to back:
        # slide.shapes is index-based. The back is the first element in the _spTree.
        # index 0 is nvSpPr, 1 is grpSpPr, then shapes.
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)
        
    prs.save(output_path)
    print(f"Saved modified PPTX with extracted background to {output_path}")

if __name__ == "__main__":
    apply_image_background("../Documents/ima/高校教师实训营.pptx", "page_2_img_1.jpeg", "../Documents/ima/高校教师实训营_v3.pptx")
